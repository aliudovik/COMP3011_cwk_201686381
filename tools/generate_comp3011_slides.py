from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


SLIDES = [
    (
        "COMP3011 Coursework 1 — drVibey API",
        [
            "Web Services and Web Data (COMP3011)",
            "Individual API Development Project",
            "Student: Anton Liudovyk / 20168381",
        ],
    ),
    (
        "Problem and Scope",
        [
            "Build a data-driven API with DB-backed CRUD and robust error handling.",
            "Domain: listener taste profiling + tailored music generation.",
            "Coursework framing: explicit CRUD over Generation resource.",
            "Demonstrable locally via Docker Compose.",
        ],
    ),
    (
        "Architecture Overview",
        [
            "Backend: Flask 3 + SQLAlchemy",
            "Database: Postgres",
            "Async pipeline: Redis + RQ worker",
            "External services: Cerebras (profile/prompt), Suno (generation)",
            "Companion frontend for demonstration flow",
        ],
    ),
    (
        "Data Model and CRUD Mapping",
        [
            "Primary assessed model: Generation",
            "Create: POST /api/generate",
            "Read: GET /api/generation/{id}",
            "Update: PATCH /api/generation/{id}",
            "Delete: DELETE /api/generation/{id}",
            "Analytics endpoint: GET /api/analytics/generations/summary?days=30",
        ],
    ),
    (
        "API Documentation Overview",
        [
            "Deliverable PDF: docs/API_Documentation.pdf",
            "Contains endpoints, parameters, auth flow, and JSON examples",
            "Documents status/error codes and response envelope",
            "Includes end-to-end CRUD demonstration sequence",
        ],
    ),
    (
        "Testing and Error Handling",
        [
            "API tests: tests/test_api_generation_crud.py",
            "Covers CRUD lifecycle, unauthorized, and validation errors",
            "Playwright smoke tests: tests/example.spec.ts",
            "Consistent JSON success/error structure across key endpoints",
        ],
    ),
    (
        "Version Control Practices and Commit History",
        [
            "Repository maintained with iterative commits",
            "Milestones: CRUD route additions, tests, docs/report deliverables",
            "Final state reflects runnable code + assessed artifacts",
            "Commit history available for examiner inspection during Q&A",
        ],
    ),
    (
        "Technical Report Highlights",
        [
            "Deliverable PDF: docs/Technical_API_Report.pdf",
            "Explains architecture, stack choices, and design rationale",
            "Reflects on testing strategy, limitations, and future improvements",
            "Includes links and references to required submission artifacts",
        ],
    ),
    (
        "GenAI Declaration and Reflection",
        [
            "Deliverable PDF: docs/GenAI_Declaration_Appendix.pdf",
            "Declares tools, purpose, and conversation-log excerpts",
            "Explains verification process and manual quality control",
            "Demonstrates methodical and transparent GenAI usage",
        ],
    ),
    (
        "All Deliverables and Oral Demo Checklist",
        [
            "Code repository + README setup instructions",
            "API documentation PDF",
            "Technical report PDF",
            "Presentation slides (this PPTX)",
            "Oral demo: start services, run CRUD flow, show analytics, defend design",
        ],
    ),
]


def main() -> None:
    repo_root = Path(__file__).resolve().parents[1]
    output_path = repo_root / "docs" / "COMP3011_Presentation_Slides.pptx"

    prs = Presentation()

    for idx, (title, bullets) in enumerate(SLIDES):
        if idx == 0:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join(bullets)
            continue

        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(22)

    prs.save(output_path)
    print(f"Created: {output_path}")


if __name__ == "__main__":
    main()
